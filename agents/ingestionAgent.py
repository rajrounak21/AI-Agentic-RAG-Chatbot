import os
import docx
import fitz
import pptx
import pandas as pd
from langchain.text_splitter import RecursiveCharacterTextSplitter
from langchain.docstore.document import Document

class IngestionAgent:
    def __init__(self, chunk_size=500, chunk_overlap=50):
        self.splitter = RecursiveCharacterTextSplitter(chunk_size=chunk_size, chunk_overlap=chunk_overlap)

    def parse_pdf(self, path):
        doc = fitz.open(path)
        return "\n".join([page.get_text() for page in doc])

    def parse_docx(self, path):
        doc = docx.Document(path)
        return "\n".join([para.text for para in doc.paragraphs])

    def parse_txt(self, path):
        with open(path, "r", encoding="utf-8") as f:
            return f.read()

    def parse_md(self, path):
        return self.parse_txt(path)

    def parse_csv(self, path):
        df = pd.read_csv(path)
        return df.to_string(index=False)

    def parse_pptx(self, path):
        prs = pptx.Presentation(path)
        texts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    texts.append(shape.text)
        return "\n".join(texts)

    def extract_text(self, path):
        ext = path.split(".")[-1].lower()
        if ext == "pdf":
            return self.parse_pdf(path)
        elif ext == "docx":
            return self.parse_docx(path)
        elif ext == "pptx":
            return self.parse_pptx(path)
        elif ext == "csv":
            return self.parse_csv(path)
        elif ext in ["txt", "md"]:
            return self.parse_txt(path)
        else:
            raise ValueError(f"Unsupported file type: {ext}")

    def run(self, file_paths):
        all_docs = []
        for file_path in file_paths:
            print(f"[IngestionAgent] Processing: {file_path}")
            text = self.extract_text(file_path)
            file_name = os.path.basename(file_path)
            chunks = self.splitter.create_documents([text], metadatas=[{"source": file_name}])
            all_docs.extend(chunks)
        return all_docs


